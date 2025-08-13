from flask import Flask, request, jsonify, send_file, url_for
import requests
import os
import tempfile
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io
import base64
from werkzeug.utils import secure_filename
import uuid
from datetime import datetime

app = Flask(__name__)

# Configure upload settings
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size
ALLOWED_EXTENSIONS = {'pdf'}

# Create uploads directory if it doesn't exist
UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

def allowed_file(filename):
    """Check if file extension is allowed"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(pdf_content):
    """Extract text from PDF content"""
    try:
        pdf_reader = PdfReader(io.BytesIO(pdf_content))
        text_content = []
        
        for page_num, page in enumerate(pdf_reader.pages):
            text = page.extract_text()
            if text.strip():
                text_content.append({
                    'page': page_num + 1,
                    'text': text.strip()
                })
        
        return text_content
    except Exception as e:
        raise Exception(f"Error extracting text from PDF: {str(e)}")

def create_ppt_from_text(text_content, output_path):
    """Create PowerPoint from extracted text with professional formatting and no overflow"""
    try:
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "PDF to PPT Conversion"
        subtitle.text = f"Converted from PDF with {len(text_content)} pages"
        
        # Add content slides
        for item in text_content:
            # Clean and format the text
            text = item['text']
            lines = [line.strip() for line in text.split('\n') if line.strip()]
            
            # Process text for better formatting
            bullet_points = []
            current_point = ""
            
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                
                # Skip very short lines
                if len(line) < 3:
                    continue
                
                # Check if this is a new section/heading
                if (line.isupper() and len(line) < 50) or line.endswith(':') or len(line) < 30:
                    if current_point:
                        bullet_points.append(current_point)
                        current_point = ""
                    bullet_points.append(line)
                else:
                    if current_point:
                        current_point += f" {line}"
                    else:
                        current_point = line
                    
                    if len(current_point) > 80:  # Shorter points
                        bullet_points.append(current_point)
                        current_point = ""
            
            if current_point:
                bullet_points.append(current_point)
            
            # Split content into multiple slides if too long
            max_points_per_slide = 6
            slide_number = 1
            
            for i in range(0, len(bullet_points), max_points_per_slide):
                slide_points = bullet_points[i:i + max_points_per_slide]
                
                # Create slide
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                # Set title
                if len(bullet_points) > max_points_per_slide:
                    title.text = f"Page {item['page']} (Part {slide_number})"
                else:
                    title.text = f"Page {item['page']}"
                
                # Add bullet points
                for j, point in enumerate(slide_points):
                    if j == 0:
                        content.text = point
                    else:
                        p = content.text_frame.add_paragraph()
                        p.text = point
                        p.level = 0
                
                # Adjust font size
                text_frame = content.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if len(run.text) > 60:
                            run.font.size = Inches(0.18)  # Even smaller for long text
                        else:
                            run.font.size = Inches(0.22)  # Normal size
                
                slide_number += 1
        
        # Save the presentation
        prs.save(output_path)
        return True
    except Exception as e:
        raise Exception(f"Error creating PowerPoint: {str(e)}")

@app.route('/convert', methods=['POST'])
def convert_pdf_to_ppt():
    """API endpoint to convert PDF file to PPT and return URL"""
    try:
        # Check if file is present in request
        if 'pdf_file' not in request.files:
            return jsonify({'error': 'No PDF file provided. Please upload a PDF file.'}), 400
        
        file = request.files['pdf_file']
        
        # Check if file is selected
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        # Check if file type is allowed
        if not allowed_file(file.filename):
            return jsonify({'error': 'Only PDF files are allowed'}), 400
        
        # Read the uploaded file
        pdf_content = file.read()
        
        if not pdf_content:
            return jsonify({'error': 'Empty file uploaded'}), 400
        
        # Extract text
        text_content = extract_text_from_pdf(pdf_content)
        
        if not text_content:
            return jsonify({'error': 'No text content found in PDF'}), 400
        
        # Generate unique filename
        original_filename = secure_filename(file.filename)
        base_filename = os.path.splitext(original_filename)[0]
        unique_id = str(uuid.uuid4())[:8]
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        ppt_filename = f"{base_filename}_{timestamp}_{unique_id}.pptx"
        ppt_path = os.path.join(app.config['UPLOAD_FOLDER'], ppt_filename)
        
        # Create PowerPoint
        create_ppt_from_text(text_content, ppt_path)
        
        # Generate download URL
        download_url = url_for('download_ppt', filename=ppt_filename, _external=True)
        
        # Return success response with URL in requested order
        response_data = {
            'success': True,
            'message': 'PDF successfully converted to PowerPoint',
            'download_url': download_url,
            'filename': ppt_filename,
            'original_filename': original_filename,
            'pages_converted': len(text_content),
            'conversion_time': datetime.now().isoformat()
        }
        return jsonify(response_data)
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/download/<filename>')
def download_ppt(filename):
    """Download converted PowerPoint file"""
    try:
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        if os.path.exists(file_path):
            return send_file(
                file_path,
                as_attachment=True,
                download_name=filename,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
        else:
            return jsonify({'error': 'File not found'}), 404
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/convert-url', methods=['POST'])
def convert_pdf_url_to_ppt():
    """API endpoint to convert PDF from URL to PPT and return URL"""
    try:
        data = request.get_json()
        
        if not data or 'pdf_url' not in data:
            return jsonify({'error': 'PDF URL is required'}), 400
        
        pdf_url = data['pdf_url']
        
        # Download PDF
        response = requests.get(pdf_url, timeout=30)
        response.raise_for_status()
        pdf_content = response.content
        
        # Extract text
        text_content = extract_text_from_pdf(pdf_content)
        
        if not text_content:
            return jsonify({'error': 'No text content found in PDF'}), 400
        
        # Generate unique filename
        unique_id = str(uuid.uuid4())[:8]
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        ppt_filename = f"converted_{timestamp}_{unique_id}.pptx"
        ppt_path = os.path.join(app.config['UPLOAD_FOLDER'], ppt_filename)
        
        # Create PowerPoint
        create_ppt_from_text(text_content, ppt_path)
        
        # Generate download URL
        download_url = url_for('download_ppt', filename=ppt_filename, _external=True)
        
        # Return success response with URL in requested order
        response_data = {
            'success': True,
            'message': 'PDF successfully converted to PowerPoint',
            'download_url': download_url,
            'filename': ppt_filename,
            'source_url': pdf_url,
            'pages_converted': len(text_content),
            'conversion_time': datetime.now().isoformat()
        }
        return jsonify(response_data)
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({'status': 'healthy', 'message': 'PDF to PPT API is running'})

@app.route('/', methods=['GET'])
def home():
    """Home endpoint with usage instructions"""
    return jsonify({
        'message': 'PDF to PPT Conversion API',
        'endpoints': {
            'convert_file': {
                'endpoint': '/convert',
                'method': 'POST',
                'type': 'multipart/form-data',
                'body': {
                    'pdf_file': 'PDF file to upload'
                },
                'response': 'JSON with download URL'
            },
            'convert_url': {
                'endpoint': '/convert-url',
                'method': 'POST',
                'type': 'application/json',
                'body': {
                    'pdf_url': 'URL of the PDF to convert'
                },
                'response': 'JSON with download URL'
            },
            'download': {
                'endpoint': '/download/<filename>',
                'method': 'GET',
                'response': 'PowerPoint file (.pptx)'
            }
        },
        'file_requirements': {
            'max_size': '16MB',
            'format': 'PDF only'
        }
    })

@app.route('/upload', methods=['GET'])
def upload_form():
    """Serve the HTML upload form"""
    return send_file('upload_form.html')

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000) 